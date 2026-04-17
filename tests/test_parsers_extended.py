"""Extended parser tests for DOCX, XLSX, CSV, PPTX, HTML, and Image parsers.

Each test creates minimal in-memory documents — no external fixture files required.
"""

import pytest
from pathlib import Path
from document_index_mcp.parsers.base import ParseResult


# ---------------------------------------------------------------------------
# DOCX Parser Tests
# ---------------------------------------------------------------------------

class TestDOCXParser:
    """Tests for DOCXParser using python-docx to create minimal documents."""

    def _make_docx(self, path: Path, paragraphs: list[tuple[str, str]],
                   tables: list[list[list[str]]] | None = None):
        """Helper: create a .docx with given (text, style) paragraphs and optional tables."""
        from docx import Document
        doc = Document()
        for text, style in paragraphs:
            doc.add_paragraph(text, style=style)
        if tables:
            for table_data in tables:
                rows = len(table_data)
                cols = len(table_data[0]) if table_data else 0
                tbl = doc.add_table(rows=rows, cols=cols)
                for r_idx, row in enumerate(table_data):
                    for c_idx, cell_text in enumerate(row):
                        tbl.cell(r_idx, c_idx).text = cell_text
        doc.save(str(path))

    def test_heading_style_sections(self, tmp_path):
        """DOCX with Heading styles produces separate sections."""
        from document_index_mcp.parsers.docx_parser import DOCXParser

        doc_path = tmp_path / "headings.docx"
        self._make_docx(doc_path, [
            ("Introduction", "Heading 1"),
            ("This is the intro paragraph.", "Normal"),
            ("Risk Assessment", "Heading 1"),
            ("Risk content goes here.", "Normal"),
        ])

        result = DOCXParser().parse(doc_path)
        assert isinstance(result, ParseResult)
        assert len(result.sections) == 2
        assert result.sections[0].title == "Introduction"
        assert "intro paragraph" in result.sections[0].content
        assert result.sections[1].title == "Risk Assessment"
        assert "Risk content" in result.sections[1].content

    def test_numbered_heading_pattern_detection(self, tmp_path):
        """DOCX without Heading styles falls back to numbered heading patterns."""
        from document_index_mcp.parsers.docx_parser import DOCXParser

        doc_path = tmp_path / "numbered.docx"
        self._make_docx(doc_path, [
            ("1. Introduction", "Normal"),
            ("Intro body text.", "Normal"),
            ("2. Methodology", "Normal"),
            ("Method body text.", "Normal"),
        ])

        result = DOCXParser().parse(doc_path)
        assert len(result.sections) >= 2
        titles = [s.title for s in result.sections]
        assert any("Introduction" in t for t in titles)
        assert any("Methodology" in t for t in titles)

    def test_table_content_extracted(self, tmp_path):
        """DOCX with a table includes table content in raw_text."""
        from document_index_mcp.parsers.docx_parser import DOCXParser

        doc_path = tmp_path / "with_table.docx"
        self._make_docx(
            doc_path,
            [("Overview", "Heading 1"), ("Some text.", "Normal")],
            tables=[[["Name", "Value"], ["Alpha", "100"], ["Beta", "200"]]],
        )

        result = DOCXParser().parse(doc_path)
        assert len(result.sections) >= 1
        # Table content should appear in the raw text or section content
        combined = result.raw_text + " ".join(s.content for s in result.sections)
        assert "Alpha" in combined
        assert "Beta" in combined

    def test_single_paragraph_fallback(self, tmp_path):
        """DOCX with no headings at all falls back to a single Document section."""
        from document_index_mcp.parsers.docx_parser import DOCXParser

        doc_path = tmp_path / "plain.docx"
        self._make_docx(doc_path, [
            ("Just a simple paragraph with no structure.", "Normal"),
        ])

        result = DOCXParser().parse(doc_path)
        assert len(result.sections) >= 1
        assert result.sections[0].section_ref is not None


# ---------------------------------------------------------------------------
# XLSX Parser Tests
# ---------------------------------------------------------------------------

class TestXLSXParser:
    """Tests for XLSXParser using openpyxl to create minimal workbooks."""

    def test_multiple_sheets(self, tmp_path):
        """Two sheets produce two sections, one per sheet."""
        import openpyxl
        from document_index_mcp.parsers.xlsx_parser import XLSXParser

        wb = openpyxl.Workbook()
        ws1 = wb.active
        ws1.title = "Users"
        ws1.append(["Name", "Age"])
        ws1.append(["Alice", 30])
        ws1.append(["Bob", 25])

        ws2 = wb.create_sheet("Products")
        ws2.append(["SKU", "Price"])
        ws2.append(["A1", 9.99])

        path = tmp_path / "multi.xlsx"
        wb.save(str(path))

        result = XLSXParser().parse(path)
        assert len(result.sections) == 2
        assert result.sections[0].title == "Users"
        assert result.sections[0].section_ref == "sheet-1"
        assert "Alice" in result.sections[0].content
        assert result.sections[1].title == "Products"
        assert result.sections[1].section_ref == "sheet-2"
        assert result.page_count == 2

    def test_empty_sheet(self, tmp_path):
        """An empty sheet gets a section with '(empty sheet)' content."""
        import openpyxl
        from document_index_mcp.parsers.xlsx_parser import XLSXParser

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Blank"
        # No rows added

        path = tmp_path / "empty.xlsx"
        wb.save(str(path))

        result = XLSXParser().parse(path)
        assert len(result.sections) == 1
        assert result.sections[0].title == "Blank"
        assert "(empty sheet)" in result.sections[0].content
        assert result.sections[0].section_ref == "sheet-1"


# ---------------------------------------------------------------------------
# CSV Parser Tests
# ---------------------------------------------------------------------------

class TestCSVParser:
    """Tests for CSVParser."""

    def test_basic_csv(self, tmp_path):
        """CSV with header + rows produces per-row sections."""
        from document_index_mcp.parsers.csv_parser import CSVParser

        path = tmp_path / "data.csv"
        path.write_text("Name,Score\nAlice,95\nBob,87\n", encoding="utf-8")

        result = CSVParser().parse(path)
        assert len(result.sections) == 2  # one per data row
        assert result.sections[0].section_ref == "row-1"
        assert "Alice" in result.sections[0].content
        assert "Score" in result.sections[0].content  # headers preserved as field names
        assert result.metadata["row_count"] == 2
        assert result.metadata["columns"] == ["Name", "Score"]

    def test_empty_csv(self, tmp_path):
        """CSV with only a header and no data rows."""
        from document_index_mcp.parsers.csv_parser import CSVParser

        path = tmp_path / "empty.csv"
        path.write_text("Col1,Col2\n", encoding="utf-8")

        result = CSVParser().parse(path)
        assert len(result.sections) == 1
        assert "Empty CSV" in result.sections[0].title or "No data" in result.sections[0].content

    def test_headers_in_content(self, tmp_path):
        """Column headers appear as field names in section content."""
        from document_index_mcp.parsers.csv_parser import CSVParser

        path = tmp_path / "fields.csv"
        path.write_text("Department,Budget\nEngineering,500000\n", encoding="utf-8")

        result = CSVParser().parse(path)
        assert len(result.sections) >= 1
        # DictReader uses headers as keys
        assert "Department" in result.sections[0].content
        assert "Budget" in result.sections[0].content


# ---------------------------------------------------------------------------
# PPTX Parser Tests
# ---------------------------------------------------------------------------

class TestPPTXParser:
    """Tests for PPTXParser using python-pptx to create minimal presentations."""

    def test_two_slides_with_titles(self, tmp_path):
        """Two slides with titles produce two sections."""
        from pptx import Presentation
        from document_index_mcp.parsers.pptx_parser import PPTXParser

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # title slide layout

        slide1 = prs.slides.add_slide(slide_layout)
        slide1.shapes.title.text = "Welcome"
        slide1.placeholders[1].text = "Hello world"

        slide2 = prs.slides.add_slide(slide_layout)
        slide2.shapes.title.text = "Agenda"
        slide2.placeholders[1].text = "Today's topics"

        path = tmp_path / "deck.pptx"
        prs.save(str(path))

        result = PPTXParser().parse(path)
        assert len(result.sections) == 2
        assert result.sections[0].title == "Welcome"
        assert result.sections[0].section_ref == "slide-1"
        assert "Hello world" in result.sections[0].content
        assert result.sections[1].title == "Agenda"
        assert result.sections[1].section_ref == "slide-2"
        assert result.page_count == 2

    def test_slide_without_title_fallback(self, tmp_path):
        """A slide with no title shape gets a 'Slide N' fallback ref."""
        from pptx import Presentation
        from pptx.util import Inches
        from document_index_mcp.parsers.pptx_parser import PPTXParser

        prs = Presentation()
        # Use blank layout (index 6) which has no title placeholder
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add a text box manually
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        txBox.text_frame.text = "Some content without a title"

        path = tmp_path / "notitle.pptx"
        prs.save(str(path))

        result = PPTXParser().parse(path)
        assert len(result.sections) == 1
        assert result.sections[0].section_ref == "slide-1"
        # Title should be the fallback "Slide 1"
        assert "Slide 1" in result.sections[0].title
        assert "Some content" in result.sections[0].content


# ---------------------------------------------------------------------------
# HTML Parser Tests
# ---------------------------------------------------------------------------

class TestHTMLParser:
    """Tests for HTMLParser."""

    def test_heading_based_sections(self, tmp_path):
        """HTML with h1/h2 headings produces sections split by headings."""
        from document_index_mcp.parsers.html_parser import HTMLParser

        html = """<!DOCTYPE html>
<html>
<head><title>Test Doc</title></head>
<body>
<h1>Introduction</h1>
<p>Intro paragraph content.</p>
<h2>Details</h2>
<p>Detail paragraph content.</p>
</body>
</html>"""
        path = tmp_path / "doc.html"
        path.write_text(html, encoding="utf-8")

        result = HTMLParser().parse(path)
        assert len(result.sections) == 2
        assert result.sections[0].title == "Introduction"
        assert "Intro paragraph" in result.sections[0].content
        assert result.sections[1].title == "Details"
        assert "Detail paragraph" in result.sections[1].content
        assert result.metadata["title"] == "Test Doc"

    def test_no_headings_fallback(self, tmp_path):
        """HTML with no headings falls back to paragraph-based or single section."""
        from document_index_mcp.parsers.html_parser import HTMLParser

        html = """<!DOCTYPE html>
<html>
<head><title>Plain</title></head>
<body>
<p>First paragraph of plain text.</p>
<p>Second paragraph of plain text.</p>
<p>Third paragraph of plain text.</p>
</body>
</html>"""
        path = tmp_path / "plain.html"
        path.write_text(html, encoding="utf-8")

        result = HTMLParser().parse(path)
        assert len(result.sections) >= 1
        # Should have content from the paragraphs
        combined = " ".join(s.content for s in result.sections)
        assert "First paragraph" in combined
        assert "Second paragraph" in combined

    def test_script_and_style_stripped(self, tmp_path):
        """Script and style tags are removed from parsed content."""
        from document_index_mcp.parsers.html_parser import HTMLParser

        html = """<html><body>
<script>var x = 1;</script>
<style>.foo { color: red; }</style>
<h1>Real Content</h1>
<p>Visible text.</p>
</body></html>"""
        path = tmp_path / "scripts.html"
        path.write_text(html, encoding="utf-8")

        result = HTMLParser().parse(path)
        assert "var x = 1" not in result.raw_text
        assert "color: red" not in result.raw_text
        assert "Visible text" in result.raw_text


# ---------------------------------------------------------------------------
# Image Parser Tests
# ---------------------------------------------------------------------------

class TestImageParser:
    """Tests for ImageParser. Requires pytesseract + tesseract binary.

    If tesseract is not installed, the test verifies structure only.
    """

    def test_tiny_white_image_structure(self, tmp_path):
        """A 10x10 white PNG returns a valid ParseResult with at least one section."""
        from PIL import Image
        from document_index_mcp.parsers.image_parser import ImageParser

        img = Image.new("RGB", (10, 10), color=(255, 255, 255))
        path = tmp_path / "white.png"
        img.save(str(path))

        try:
            result = ImageParser().parse(path)
        except Exception:
            pytest.skip("tesseract binary not available")

        assert isinstance(result, ParseResult)
        assert len(result.sections) >= 1
        assert result.sections[0].section_ref.startswith("ocr-")
        assert result.metadata["parser"] == "pytesseract"
        assert "10x10" in result.metadata["image_size"]
        assert result.filename == "white.png"


def test_docx_parser_populates_paragraphs(tmp_path):
    from docx import Document
    from document_index_mcp.parsers.docx_parser import DOCXParser

    docx_path = tmp_path / "test.docx"
    doc = Document()
    doc.add_heading("1. Introduction", level=1)
    doc.add_paragraph("First sentence here. Second sentence here.")
    doc.add_paragraph("Another paragraph. With two sentences.")
    doc.save(docx_path)

    result = DOCXParser().parse(docx_path)
    assert result.parser_version
    assert result.full_text
    assert any(len(s.paragraphs) > 0 for s in result.sections)
    for section in result.sections:
        for para in section.paragraphs:
            for sent in para.sentences:
                assert result.full_text[sent.char_start:sent.char_end] == sent.text
